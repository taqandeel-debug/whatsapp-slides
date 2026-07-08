import streamlit as st
from pptx import Presentation
from pptx.util import Inches
import io

# Page Config optimized for mobile viewports
st.set_page_config(page_title="WhatsApp Slide Maker", layout="wide")

st.title("📱 Daily WhatsApp Slides")
st.write("Upload photos. They will stay safely in memory until you download.")

# Initialize background memory (Session State) so photos never disappear
if "saved_photos" not in st.session_state:
    st.session_state.saved_photos = []

# 1. File Uploader
uploaded_files = st.file_uploader(
    "Tap to add photos", 
    accept_multiple_files=True, 
    type=['png', 'jpg', 'jpeg'],
    key="uploader_input"
)

# If new files are uploaded, instantly save them into background memory
if uploaded_files:
    # Merge or overwrite background memory with new files
    st.session_state.saved_photos = uploaded_files
    # Chronological sort based on WhatsApp file names (IMG-YYYYMMDD...)
    st.session_state.saved_photos.sort(key=lambda x: x.name)

# Process only if we have photos safely tucked in background memory
if st.session_state.saved_photos:
    total_photos = len(st.session_state.saved_photos)
    st.success(f"✅ {total_photos} photos safely holding in memory.")
    
    # 2. Preview Layout
    st.divider()
    st.subheader("👀 Slide Preview")
    
    # Group images into blocks of 3
    chunks = [st.session_state.saved_photos[i:i + 3] for i in range(0, total_photos, 3)]
    
    for i, group in enumerate(chunks):
        st.write(f"**Slide {i+1}**")
        cols = st.columns(3)
        for idx, file in enumerate(group):
            # Ensure file pointer is at the start before viewing
            file.seek(0)
            cols[idx].image(file, use_container_width=True)
            cols[idx].caption(f"Pic {idx+1}")
    
    st.divider()

    # 3. Secure PowerPoint Generation (No memory loss on click)
    # Using container to avoid full page refresh issues on mobile browsers
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6] # Secure standard blank layout index
    
    for group in chunks:
        slide = prs.slides.add_slide(blank_layout)
        positions = [Inches(0.5), Inches(4.7), Inches(8.9)]
        top = Inches(1.5)
        height = Inches(4.5)
        
        for idx, file in enumerate(group):
            file.seek(0) # Reset pointer safely
            slide.shapes.add_picture(file, positions[idx], top, height=height)

    # Save presentation to memory buffer
    ppt_output = io.BytesIO()
    prs.save(ppt_output)
    ppt_output.seek(0)
    
    # Direct download button (Eliminates the unstable two-button step)
    st.download_button(
        label="📥 Download PowerPoint File",
        data=ppt_output,
        file_name="Daily_Updates.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        use_container_width=True
    )
    
    # Reset button to clear memory for next batch
    if st.button("🗑️ Clear All Photos & Start New", use_container_width=True):
        st.session_state.saved_photos = []
        st.rerun()
